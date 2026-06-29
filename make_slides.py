from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ── palette ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x35, 0x64)   # slide background / title bar
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xF0, 0xB3, 0x23)   # accent stripe / headings
LIGHT  = RGBColor(0xF2, 0xF4, 0xF8)   # body background on body slides
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x1A, 0x8C, 0x4E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── helpers ──────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 5.8, 13.33, 0.25, GOLD)
    add_text(slide, title,   1, 2.2, 11.33, 1.5, font_size=40, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, 1, 3.9, 11.33, 0.8, font_size=22,
                 color=GOLD, align=PP_ALIGN.CENTER)


def section_slide(title):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0.5, 3.3, 12.33, 0.08, GOLD)
    add_text(slide, title, 0.5, 2.5, 12.33, 1.5, font_size=34, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)


def body_slide(title, bullets, icons=None):
    """
    bullets: list of (text, level, flag)
      level 0 = main bullet  level 1 = sub-bullet
      flag: 'ok' | 'warn' | 'bad' | None
    icons: optional dict mapping flag → emoji prefix (unused if None)
    """
    slide = prs.slides.add_slide(BLANK)
    # background
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT)
    # title bar
    add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
    add_rect(slide, 0, 1.1, 13.33, 0.06, GOLD)
    add_text(slide, title, 0.3, 0.12, 12.5, 0.9, font_size=26, bold=True,
             color=WHITE)

    # bullet area
    txb = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.9)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for (text, level, flag) in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.level = level
        indent = Inches(0.3 * level)
        p.space_before = Pt(4 if level == 0 else 2)

        # bullet symbol
        run = p.add_run()
        if level == 0:
            run.text = "▸  " + text
        else:
            run.text = "    • " + text

        sz = 17 if level == 0 else 15
        run.font.size = Pt(sz)
        run.font.bold = (level == 0)

        if flag == 'bad':
            run.font.color.rgb = RED
        elif flag == 'ok':
            run.font.color.rgb = GREEN
        else:
            run.font.color.rgb = NAVY

    return slide


# ══════════════════════════════════════════════════════════════════════════
# SLIDE CONTENT
# ══════════════════════════════════════════════════════════════════════════

# 1. Title
title_slide("May 2026 Clinical Walkthrough",
            "Multi-Site Quality Review  |  May 2026")

# 2. Agenda
body_slide("Overview", [
    ("Sites Visited",         0, None),
    ("Princeton  –  May 6",   1, None),
    ("Bethesda  –  May 13",   1, None),
    ("Bowie  –  May 14",      1, None),
    ("Key Themes & Action Items", 0, None),
])

# ── PRINCETON ─────────────────────────────────────────────────────────────
section_slide("Princeton  |  May 6, 2026")

body_slide("Princeton – Clinic Overview", [
    ("20 patients  •  16 scans  •  3 techs",   0, None),
    ("Issues & Complaints",                     0, 'bad'),
    ("Medications not labeled",                 1, 'bad'),
    ("Storage room could use an additional shelf", 1, 'warn'),
])

body_slide("Princeton – Varithena (Cortney)", [
    ("20 minutes total",                           0, None),
    ("Measured all artifact",                      1, None),
    ("DVT r/o performed before procedure",         1, 'ok'),
    ("Optisign video played",                      1, 'ok'),
    ("Consents signed",                            1, 'ok'),
    ("Sharps disposed of properly",                1, 'ok'),
    ("Time-out performed out loud",                1, 'ok'),
    ("Areas for Improvement",                      0, 'warn'),
    ("Scanned SSV but documented as GSV",          1, 'bad'),
    ('Patient felt experience focused on cosmetics — not her concern', 1, 'warn'),
])

body_slide("Princeton – Varithena (McKinley)", [
    ("8 minutes total",                 0, None),
    ("Talk ran 12 minutes",             1, 'warn'),
    ("Presenter not fully confident",   1, 'warn'),
    ('"Knock on wood" language used',   1, 'warn'),
    ("1 hour total appointment",        1, None),
    ("Consents signed",                 1, 'ok'),
    ("Sharps disposed of properly",     1, 'ok'),
    ("Time-out performed out loud",     1, 'ok'),
])

body_slide("Princeton – RFA (Jen)", [
    ("23 minutes total  (13 min for RFA)",         0, None),
    ("Scan took 40 minutes  –  1 hr total",        1, 'warn'),
    ("Optisign video NOT played",                  1, 'bad'),
    ("Initially refused scan-assist; accepted after trial", 1, 'warn'),
    ("Consent signed",                             1, 'ok'),
    ("Sharps disposed of properly",               1, 'ok'),
    ("Time-out performed out loud",               1, 'ok'),
    ("Initials: Jen",                             1, None),
])

# ── BETHESDA ──────────────────────────────────────────────────────────────
section_slide("Bethesda  |  May 13, 2026")

body_slide("Bethesda – Clinic Overview", [
    ("26 patients  •  18 scans  •  3 techs",   0, None),
    ("Issues & Complaints",                     0, 'bad'),
    ("Front desk rooms all patients",           1, 'warn'),
    ("Techs prep RFAs independently",           1, 'warn'),
    ("Medications not labeled",                 1, 'bad'),
])

body_slide("Bethesda – Varithena + US (Rob)", [
    ("24 minutes total  (8 min for Varithena)",    0, None),
    ("All consents signed",                        1, 'ok'),
    ("Varithena tray set-up was efficient",        1, 'ok'),
    ("Rob was very outgoing with patient",         1, 'ok'),
    ("Sharps disposed of properly",               1, 'ok'),
    ("Time-out NOT performed out loud",            1, 'bad'),
])

body_slide("Bethesda – US Review", [
    ("14 minutes total",                       0, None),
    ("Thoroughly explained findings to patient", 1, 'ok'),
    ("Folder given by front desk",             1, 'ok'),
    ("Optisign video played",                  1, 'ok'),
    ("Note completed properly",                1, 'ok'),
    ("No iPad available",                      1, 'bad'),
])

body_slide("Bethesda – RFA + US (Dr. Green / Ruth – Training)", [
    ("40 minutes total  (14 min for RFA)",         0, None),
    ("Consent signed",                             1, 'ok'),
    ("Sharps disposed of properly",               1, 'ok'),
    ("Note completed",                            1, 'ok'),
    ("Time-out NOT performed out loud",            1, 'bad'),
    ("Wrong note opened initially",               1, 'bad'),
    ("US note not completed properly",            1, 'bad'),
])

# ── BOWIE ──────────────────────────────────────────────────────────────────
section_slide("Bowie  |  May 14, 2026")

body_slide("Bowie – Clinic Overview", [
    ("Issues & Complaints",                     0, 'bad'),
    ("Front desk rooms all patients",           1, 'warn'),
    ("RFA times are stretched",                 1, 'warn'),
    ("Techs prep RFAs independently",           1, 'warn'),
])

body_slide("Bowie – RFA + US (Case 1)", [
    ("55 minutes total  (15 min for RFA)",             0, None),
    ("Sharps disposed of properly",                    1, 'ok'),
    ("Topical lidocaine used — vein smaller, no access issues", 1, 'ok'),
    ("Time-out NOT performed out loud",                1, 'bad'),
    ("16g needle used before intro noted",             1, 'bad'),
    ("Physician bent over during ablation",            1, 'warn'),
    ("Sonographer assisting patient cleaning — limits availability for next appt", 1, 'warn'),
])

body_slide("Bowie – RFA + US (Case 2)", [
    ("1.5 hours total  (RFA 35 min)",                  0, None),
    ("Sharps disposed of properly",                    1, 'ok'),
    ("Time-out NOT performed out loud",                1, 'bad'),
    ("RFA performed before post-scan",                 1, 'bad'),
    ("Multiple failed access attempts on GSV  →  switched to SSV", 1, 'warn'),
    ("Physician repeatedly looked to observer for guidance", 1, 'warn'),
    ("US note not completed properly",                 1, 'bad'),
])

body_slide("Bowie – US Review", [
    ("10 minutes total",                           0, None),
    ("Thoroughly explained findings to patient",   1, 'ok'),
    ("Folder given by front desk",                 1, 'ok'),
    ("Optisign video played",                      1, 'ok'),
    ("Note completed",                             1, 'ok'),
    ("No iPad available",                          1, 'bad'),
])

# ── SUMMARY ────────────────────────────────────────────────────────────────
section_slide("Key Themes & Action Items")

body_slide("Recurring Issues Across All Sites", [
    ("🔴  Critical — Address Immediately",         0, 'bad'),
    ("Medications NOT labeled at all sites",       1, 'bad'),
    ("Time-out NOT performed out loud (multiple sites)", 1, 'bad'),
    ("🟡  Operational Concerns",                   0, 'warn'),
    ("Front desk rooming patients (Bethesda, Bowie)", 1, 'warn'),
    ("Techs prepping RFAs independently (Bethesda, Bowie)", 1, 'warn'),
    ("No iPad for US reviews (Bethesda, Bowie)",   1, 'warn'),
    ("Optisign video not played in some cases",    1, 'warn'),
    ("🟢  Positive Observations",                  0, 'ok'),
    ("Sharps disposal consistently compliant",     1, 'ok'),
    ("Most consents signed",                       1, 'ok'),
    ("Patient education thorough in US reviews",   1, 'ok'),
])

body_slide("Action Items", [
    ("Label ALL medications — no exceptions",                           0, 'bad'),
    ("Reinforce time-out protocol — must be performed out loud",        0, 'bad'),
    ("Clarify front desk vs. tech rooming responsibilities",            0, 'warn'),
    ("Ensure Optisign video is played for every eligible patient",      0, 'warn'),
    ("Provide iPads at Bethesda and Bowie for US reviews",              0, 'warn'),
    ("Address RFA documentation workflow — correct note, US notes complete", 0, 'warn'),
    ("Coach McKinley on procedure talk confidence and timing",          0, 'warn'),
    ("Bowie: review RFA sequencing (post-scan should follow RFA)",      0, 'warn'),
])

# save
out = "/root/.openclaw/workspace/May_2026_Walkthrough_Slides.pptx"
prs.save(out)
print(f"Saved: {out}")
